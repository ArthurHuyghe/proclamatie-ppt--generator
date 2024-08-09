# Import necessary libraries
import json  # For loading JSON files
import os  # For file path operations
import pptx  # For creating and editing PowerPoint presentations

# Main

# Configure paths
# Path to a JSON file containing the class lists
path_klaslijst = r".\example\klaslijst.json"

# Path to the folders containing the students pictures
path_fotos_1ste = r".\example\1ste jaar"
path_fotos_3de = r".\example\3de jaar"
path_fotos_5de = r".\example\5de jaar"
path_fotos_6de = r".\example\6de jaar"

# Load JSON data
with open(path_klaslijst, "r", encoding="utf-8") as file:
    klaslijsten = json.load(file)

# Open the PowerPoint presentation template
prs = pptx.Presentation("Proclamatie sjabloon.pptx")

# Add the Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.placeholders[10].text = "6de jaar 2023-2024"  # Change this to the correct date

# Start adding slides for each class and student
# Initialize an empty list to store names of students without corresponding photos
leerlingen_without_babyface = []
leerlingen_without_grown_upface = []

# Iterate through each titular (teacher) in the JSON data
for titularis in klaslijsten["titularissen"]:
    titularis_naam = titularis["naam"]  # Get the teacher's name
    for klas in titularis["klassen"]:
        klas_name = klas["klas"]  # Get the class name

        # Add a slide with the class name
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.placeholders[0].text = klas_name  # Set the class name on the slide
        slide.placeholders[10].text = (
            titularis_naam  # Set the teacher's name on the slide
        )

        # Iterate through each student in the class
        for leerling in klas["leerlingen"]:
            leerling_name = leerling["Naam"]  # Get the student's name
            # Construct file paths for the student's babyface and grown-up face photos
            babyface_path = os.path.join(
                path_fotos_1ste, f"{leerling_name.lower().replace(" ", ".")}.jpg"
            )
            grown_upface_path = os.path.join(
                path_fotos_6de, f"{leerling_name.lower().replace(" ", ".")}.jpg"
            )

            # Add a slide for the student
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            slide.placeholders[0].text = (
                leerling_name  # Set the student's name on the slide
            )
            slide.placeholders[12].text = klas_name  # Set the class name on the slide

            # Adding images
            # Try adding the babyface image
            try:
                slide.placeholders[10].insert_picture(babyface_path)
            # If the 1st-year photo is not found, try the 3rd-year photo
            except FileNotFoundError:
                babyface_path = os.path.join(
                    path_fotos_3de, f"{leerling_name.lower().replace(" ", ".")}.jpg"
                )
                try:
                    slide.placeholders[10].insert_picture(babyface_path)
                # If the 3rd-year photo is not found, try the 5th-year photo
                except FileNotFoundError:
                    babyface_path = os.path.join(
                        path_fotos_5de, f"{leerling_name.lower().replace(" ", ".")}.jpg"
                    )
                    try:
                        slide.placeholders[10].insert_picture(babyface_path)
                    # If no babyface photo is found, add the student's name to the list
                    except FileNotFoundError:
                        print(f"No Babyface found for: {leerling_name}")
                        leerlingen_without_babyface.append(leerling_name)

            # Try adding the grown-up face image
            try:
                slide.placeholders[11].insert_picture(grown_upface_path)
            # If the 6th-year photo is not found, try the 5th-year photo
            except FileNotFoundError:
                grown_upface_path = os.path.join(
                    path_fotos_5de, f"{leerling_name.lower().replace(" ", ".")}.jpg"
                )
                try:
                    slide.placeholders[11].insert_picture(grown_upface_path)
                # If no grown-up face photo is found, add the student's name to the list
                except FileNotFoundError:
                    print(f"No Grown-upface found for: {leerling_name}")
                    leerlingen_without_grown_upface.append(leerling_name)


# Save the presentation
prs.save(".\example\Proclamatie foto's.pptx")
print("Powerpoint succesfully created.")

# Print the lists of students without babyface and grown-up face photos
print(f"{len(leerlingen_without_babyface)} People without babyface:")
print(leerlingen_without_babyface)
print("")
print(f"{len(leerlingen_without_grown_upface)} People without grown-upface:")
print(leerlingen_without_grown_upface)
