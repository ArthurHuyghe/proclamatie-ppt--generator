#!/usr/bin/python3

import pptx
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor


def replace_image_on_slide(img_shape, new_image_path: str):
    # Retrieve the part and relationship ID from the shape to be changed
    slide_part, rId = img_shape.part, img_shape._element.blip_rId

    # Add the new image to the slide and get the new relationship ID
    new_image_part, new_rid = slide_part.get_or_add_image_part(new_image_path)

    # Update the shape to reference the new image
    blip_fill = img_shape._element.blipFill
    blip = blip_fill.find(qn("a:blip"))
    blip.set(qn("r:embed"), new_rid)


def change_title(slide, text: str):
    """
    Changes the title of the slide and applies the same formatting as the template slide.
    """
    # Access the title placeholder on the slide
    title_shape = slide.shapes.title
    # Change the title text to the text thats passed
    title_shape.text = text
    # Set the font, size and position for the title
    title_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.left = Cm(3.30)
    title_shape.top = Cm(1.12)
    title_shape.width = Cm(18.63)
    title_shape.height = Cm(3.18)
    # Change the color of the title to white
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


img_grown_up = "einde.jpg"
img_babyface = "begin.png"

# open presentation
prs = pptx.Presentation("Proclamatie foto's sjabloon.pptx")

# create new image part from new image file
new_pptx_img = pptx.parts.image.Image.from_file(img_grown_up)

# obviously have to figure out what image you're actually changing... ps: start counting from 0
shape_grown_up = prs.slides[4].shapes[2]
shape_babyface = prs.slides[4].shapes[3]

replace_image_on_slide(shape_grown_up, img_grown_up)
replace_image_on_slide(shape_babyface, img_babyface)

# change the title to the according name
name = "Patrick Hans"
change_title(prs.slides[4], name)


# save it
prs.save("presentation_with_picture.pptx")
