from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ppt = Presentation("Proclamatie foto's sjabloon.pptx")
for i in range(0, 1):
    # Add a slide
    slide_layout = ppt.slide_layouts[5]  # 5 is for slide layout with title
    slide = ppt.slides.add_slide(slide_layout)

    # Add the first picture to the slide
    image_path = "begin.png"
    pic = slide.shapes.add_picture(
        image_path, left=Cm(3.30), top=Cm(5.73), width=None, height=Cm(10.00)
    )

    # Add the second picture to the slide
    image_path = "einde.jpg"
    left = Cm(14.14)
    top = Cm(5.73)
    height = Cm(10.00)
    width = None
    pic = slide.shapes.add_picture(image_path, left, top, width, height)

    # Access the title placeholder on the slide
    title_shape = slide.shapes.title

    # Change the title text to "Your Title Here"
    title_shape.text = "Homer Simpson"

    # Set the font, size and position for the title
    title_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.left = Cm(3.30)
    title_shape.top = Cm(1.12)
    title_shape.width = Cm(18.63)
    title_shape.height = Cm(3.18)
    # Change the color of the title to white
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add class to the slide with a specific font and size
    left = Cm(0.70)
    top = Cm(0.72)
    width = Cm(3.00)
    height = Cm(1.62)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "6GL"
    # Set the font properties
    text_properties = text_frame.paragraphs[0]
    text_properties.font.name = "Eras Medium ITC"
    text_properties.font.size = Pt(32)
    text_properties.font.color.rgb = RGBColor(255, 255, 255)  # White color
    # Align the text in the center
    text_properties.alignment = PP_ALIGN.CENTER

# Save the presentation
ppt.save("presentation_with_picture.pptx")
