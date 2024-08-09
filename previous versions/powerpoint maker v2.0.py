#!/usr/bin/python3

import pptx
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE


def replace_image_on_slide(img_shape, new_image_path: str):
    # Retrieve the part and relationship ID from the shape to be changed
    slide_part, rId = img_shape.part, img_shape._element.blip_rId

    # Add the new image to the slide and get the new relationship ID
    new_image_part, new_rid = slide_part.get_or_add_image_part(new_image_path)

    # Update the shape to reference the new image
    blip_fill = img_shape._element.blipFill
    blip = blip_fill.find(qn("a:blip"))
    blip.set(qn("r:embed"), new_rid)


def change_title(slide, text: str):
    """
    Changes the title of the slide and applies the same formatting as the template slide.
    """
    # Access the title placeholder on the slide
    title_shape = slide.shapes.title
    # Change the title text to the text thats passed
    title_shape.text = text
    # Set the font, size and position for the title
    title_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.left = Cm(3.30)
    title_shape.top = Cm(1.12)
    title_shape.width = Cm(18.63)
    title_shape.height = Cm(3.18)
    # Change the color of the title to white
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def change_classname(slide, text: str):
    """
    Changes the classname of the slide and applies the same formatting as the template slide.
    """
    # Access the title placeholder on the slide
    class_shape = slide.shapes[1]
    # Change the title text to the text thats passed
    class_shape.text = text
    # Set the font, size and color for the classname
    class_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    class_shape.text_frame.paragraphs[0].font.size = Pt(32)
    class_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    # Align the text in the center and make the box fit
    class_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    # class_shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # class_shape.text_frame.word_wrap = False
    # Set the position and dimensions
    class_shape.left = Cm(0.70)
    class_shape.top = Cm(0.72)
    class_shape.height = Cm(1.62)
    class_shape.width = Cm(6)


img_grown_up = r"pasfotos\6e jaar 23-24\6WEWIe2\Aaron Biebau.jpg"
img_babyface = r"pasfotos\1e jaar 18-19\1L1\Aaron Biebau.jpg"

# open presentation
prs = pptx.Presentation("Proclamatie foto's sjabloon.pptx")

# create new image part from new image file
new_pptx_img = pptx.parts.image.Image.from_file(img_grown_up)

# obviously have to figure out what image you're actually changing... ps: start counting from 0
shape_grown_up = prs.slides[3].shapes[3]
shape_babyface = prs.slides[3].shapes[2]

replace_image_on_slide(shape_grown_up, img_grown_up)
replace_image_on_slide(shape_babyface, img_babyface)

# change the title to the according name
name = "Aaron Biebau"
classname = "6WEWIe2"
change_title(prs.slides[3], name)
change_classname(prs.slides[3], classname)

# save it
prs.save("presentation_with_picture.pptx")
