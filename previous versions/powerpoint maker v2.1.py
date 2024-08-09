#!/usr/bin/python3

import json
import os
import pptx
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE


def replace_image_on_slide(img_shape, new_image_path: str):
    """Replaces the image that is passed with a new image on a powerpoint slide.

    Args:
        img_shape (pptx-shape): pptx-shape of the image you want to replace.
        new_image_path (str): Path to the image you want to add.

    Returns:
        Nothing, but prints an error-message if the image was not found.
    """
    # Check if image path is valid
    if os.path.exists(new_image_path):
        # Retrieve the part and relationship ID from the shape to be changed
        slide_part, rId = img_shape.part, img_shape._element.blip_rId

        # Add the new image to the slide and get the new relationship ID
        new_image_part, new_rid = slide_part.get_or_add_image_part(new_image_path)

        # Update the shape to reference the new image
        blip_fill = img_shape._element.blipFill
        blip = blip_fill.find(qn("a:blip"))
        blip.set(qn("r:embed"), new_rid)

        return True

    else:
        return False


def change_name(slide, text: str):
    """
    Changes the title of the slide and applies the same formatting as the template slide.
    """
    # Access the title placeholder on the slide
    title_shape = slide.shapes.title
    # Change the title text to the text thats passed
    title_shape.text = text
    # Set the font, size and position for the title
    title_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.left = Cm(3.30)
    title_shape.top = Cm(1.12)
    title_shape.width = Cm(18.63)
    title_shape.height = Cm(3.18)
    # Change the color of the title to white
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def change_classname(slide, text: str):
    """
    Changes the classname of the slide and applies the same formatting as the template slide.
    """
    # Access the title placeholder on the slide
    class_shape = slide.shapes[1]
    # Change the title text to the text thats passed
    class_shape.text = text
    # Set the font, size and color for the classname
    class_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    class_shape.text_frame.paragraphs[0].font.size = Pt(32)
    class_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    # Align the text in the center and make the box fit
    class_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    # class_shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # class_shape.text_frame.word_wrap = False
    # Set the position and dimensions
    class_shape.left = Cm(0.70)
    class_shape.top = Cm(0.72)
    class_shape.height = Cm(1.62)
    class_shape.width = Cm(6)


def change_klasslide(slide, text: str):
    """
    Changes the title of the slide and applies the same formatting as the template slide.
    """
    # Change slide layout to title slide
    slide.layout = prs.slide_layouts[0]
    # Access the title placeholder on the slide
    title_shape = slide.shapes.title
    # Change the title text to the text thats passed
    title_shape.text = text
    # Set the font, size and position for the title
    title_shape.text_frame.paragraphs[0].font.name = "Eras Medium ITC"
    title_shape.text_frame.paragraphs[0].font.size = Pt(80)
    # title_shape.left = Cm(3.30)
    # title_shape.top = Cm(1.12)
    # title_shape.width = Cm(18.63)
    # title_shape.height = Cm(3.18)
    # Change the color of the title to white
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


# Main

# Configure paths
path_klaslijst = r"C:\Users\Arthu\Documents\School\python\proclamatie\klaslijsten\klaslijsten zesde jaar.json"
path_fotos_1ste = r"C:\Users\Arthu\Documents\School\python\proclamatie\pasfotos\1e jaar 18-19\allemaal samen"
path_fotos_3de = r"C:\Users\Arthu\Documents\School\python\proclamatie\pasfotos\3e Jaar 20-21\allemaal samen"
path_fotos_5de = r"C:\Users\Arthu\Documents\School\python\proclamatie\pasfotos\5e Jaar 22-23\allemaal samen"
path_fotos_6de = r"C:\Users\Arthu\Documents\School\python\proclamatie\pasfotos\6e jaar 23-24\allemaal samen"

# Load JSON data
with open(path_klaslijst, "r", encoding="utf-8") as file:
    klaslijsten = json.load(file)

# Open presentation
prs = pptx.Presentation("Proclamatie foto's sjabloon.pptx")
# Start changing process
slide_number = 0
for klasgroep in klaslijsten["klasgroepen"]:
    klasgroep_name = klasgroep["klasgroep"]
    slide_number += 1

    change_klasslide(prs.slides[slide_number], klasgroep_name)

    # Adapt each slide to the student
    for leerling in klasgroep["leerlingen"]:
        # Set variables
        slide_number += 1
        leerling_name = leerling["Naam"]
        babyface_path = os.path.join(path_fotos_1ste, f"{leerling_name}.jpg")
        grown_upface_path = os.path.join(path_fotos_6de, f"{leerling_name}.jpg")

        # Change name & class
        change_name(prs.slides[slide_number], leerling_name)
        change_classname(prs.slides[slide_number], klasgroep_name)

        # Changing babyface image
        if not replace_image_on_slide(
            prs.slides[slide_number].shapes[2], babyface_path
        ):
            print(f"Babyface not found for '{leerling_name}' at: {babyface_path}")
            babyface_path = os.path.join(
                path_fotos_3de, f"{leerling_name.lower().replace(" ", ".")}.jpg"
            )
            if not replace_image_on_slide(
                prs.slides[slide_number].shapes[2], babyface_path
            ):
                print(f"Babyface not found for '{leerling_name}' at: {babyface_path}")
            else:
                print(f"Image found for {leerling_name}")
        # Changing grown-upface image
        replace_image_on_slide(prs.slides[slide_number].shapes[3], grown_upface_path)


# Save the presentation
prs.save("Proclamatie foto's.pptx")
print(f"Powerpoint succesfully created. {slide_number} slides were changed.")
